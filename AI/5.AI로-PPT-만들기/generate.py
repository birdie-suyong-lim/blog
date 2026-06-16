"""
AI로 PPT 만들기 — 회사 입문자용 13장 가이드 덱 생성기

같은 폴더의 아티클(`../5.AI로-PPT-만들기.md`)을 입문자 청자에 맞춰
풀어낸 슬라이드. 갈래 C(LLM + 변환) 방식의 데모이기도 함.

실행:
    python3 generate.py
출력:
    AI로-PPT-만들기-입문자가이드.pptx (같은 폴더)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ────────────────────────────────────────────────────────────────────
# 디자인 토큰
# ────────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x3A, 0x68)
BLUE = RGBColor(0x4A, 0x6F, 0xA5)
ORANGE = RGBColor(0xF0, 0x8A, 0x24)
BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT_LINE = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Apple SD Gothic Neo"  # macOS 기본. 윈도우/일반 PC에선 자동 fallback


# ────────────────────────────────────────────────────────────────────
# 헬퍼
# ────────────────────────────────────────────────────────────────────
def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # 맨 뒤로
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_round_rect(slide, x, y, w, h, fill=None, line=None, line_width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.adjustments[0] = 0.15
    shape.shadow.inherit = False
    return shape


def header(slide, title, kicker=None, page=None, total=None):
    """공통 헤더: 상단 작은 카테고리 라벨 + 제목 + 좌측 액센트 바 + 우측 페이지 번호"""
    # 좌측 액센트 바
    add_rect(slide, Inches(0.7), Inches(0.55), Inches(0.08), Inches(0.55), fill=ORANGE)
    if kicker:
        add_text(slide, Inches(0.95), Inches(0.5), Inches(8), Inches(0.35),
                 kicker, size=11, bold=True, color=ORANGE)
    add_text(slide, Inches(0.95), Inches(0.78), Inches(11), Inches(0.7),
             title, size=26, bold=True, color=NAVY)
    # 하단 얇은 구분선
    add_rect(slide, Inches(0.7), Inches(1.55), Inches(11.9), Emu(9525), fill=LIGHT_LINE)
    # 페이지 번호
    if page is not None and total is not None:
        add_text(slide, Inches(11.7), Inches(7.05), Inches(1.0), Inches(0.3),
                 f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ────────────────────────────────────────────────────────────────────
# 슬라이드 빌더
# ────────────────────────────────────────────────────────────────────
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)

    # 배경 장식 도형
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(7), Inches(7))
    deco.fill.solid()
    deco.fill.fore_color.rgb = BLUE
    deco.line.fill.background()
    deco.shadow.inherit = False

    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(4), Inches(4))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = ORANGE
    deco2.line.fill.background()
    deco2.shadow.inherit = False

    add_text(s, Inches(0.7), Inches(2.0), Inches(8), Inches(0.4),
             "AI 활용 시리즈 — PPT 편", size=14, bold=True, color=ORANGE)
    add_text(s, Inches(0.7), Inches(2.5), Inches(9), Inches(1.5),
             "AI로 PPT 만들기", size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.0), Inches(10), Inches(1),
             "회사에서 AI를 막 배우기 시작하는\n사람들을 위한 30분 가이드",
             size=22, color=WHITE)
    add_text(s, Inches(0.7), Inches(6.5), Inches(6), Inches(0.35),
             "2026.05.12   |   사내 학습 자료", size=11, color=BG)


def slide_agenda(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "오늘 다룰 것", kicker="AGENDA", page=page, total=total)

    items = [
        ("01", "왜 지금 AI로 PPT인가", "슬라이드 만드는 시간이 어디로 가는지"),
        ("02", "AI PPT의 3가지 접근법", "큰 그림 한 번에"),
        ("03", "상황별로 어떤 걸 쓰는가", "내부용/외부용/반복용"),
        ("04", "30분 빠른 시작", "Gamma로 첫 덱 만들기"),
        ("05", "좋은 프롬프트 4요소", "AI 티 줄이기"),
        ("06", "흔한 함정", "환각·보안·템플릿 충돌"),
    ]

    y = 2.0
    for num, title, desc in items:
        # 번호 박스
        add_round_rect(s, Inches(0.95), Inches(y), Inches(0.7), Inches(0.7),
                       fill=NAVY)
        add_text(s, Inches(0.95), Inches(y), Inches(0.7), Inches(0.7),
                 num, size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 본문
        add_text(s, Inches(1.85), Inches(y + 0.02), Inches(9), Inches(0.4),
                 title, size=17, bold=True, color=NAVY)
        add_text(s, Inches(1.85), Inches(y + 0.4), Inches(9), Inches(0.35),
                 desc, size=12, color=MUTED)
        y += 0.85


def slide_problem(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "슬라이드 만드는 데 시간이 어디로 가나요?",
           kicker="01. 문제 의식", page=page, total=total)

    # 좌측: 큰 질문
    add_text(s, Inches(0.95), Inches(2.0), Inches(6), Inches(0.6),
             "발표 준비할 때 가장 오래 걸리는 건…", size=18, color=MUTED)
    add_text(s, Inches(0.95), Inches(2.6), Inches(6.5), Inches(2.2),
             "내용 정리가 아니라\n레이아웃·정렬·아이콘 찾기.",
             size=30, bold=True, color=NAVY)

    add_text(s, Inches(0.95), Inches(5.0), Inches(6.5), Inches(2),
             "메시지 정리에 30분,\n슬라이드로 옮기는 데 3시간.\n\n비율이 거꾸로다.",
             size=16, color=TEXT)

    # 우측: 시간 분배 비주얼
    add_text(s, Inches(8), Inches(2.0), Inches(4.5), Inches(0.4),
             "기존 방식의 시간 분배", size=12, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER)

    # 메시지 정리 (작은 박스)
    add_rect(s, Inches(8.3), Inches(2.5), Inches(0.6), Inches(0.6), fill=BLUE)
    add_text(s, Inches(8.3), Inches(2.5), Inches(0.6), Inches(0.6),
             "30m", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.0), Inches(2.65), Inches(3), Inches(0.4),
             "메시지 정리", size=12, color=TEXT)

    # 슬라이드 디자인 (큰 박스)
    add_rect(s, Inches(8.3), Inches(3.3), Inches(3.6), Inches(0.6), fill=ORANGE)
    add_text(s, Inches(8.3), Inches(3.3), Inches(3.6), Inches(0.6),
             "180m  ←  여기가 문제", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.3), Inches(4.0), Inches(3.6), Inches(0.4),
             "슬라이드 디자인·정렬·아이콘",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)


def slide_three_approaches(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "AI로 PPT 만드는 3가지 접근법",
           kicker="02. 큰 그림", page=page, total=total)

    cards = [
        ("A", "AI 네이티브 도구", "처음부터 'AI가 슬라이드를 만들도록'\n설계된 도구",
         "Gamma · Tome · Beautiful.ai", BLUE),
        ("B", "기존 도구 + AI", "PowerPoint · Google Slides · Canva\n안에 AI가 들어간 형태",
         "Copilot · Gemini · Magic Design", NAVY),
        ("C", "LLM + 변환 코드", "ChatGPT/Claude로 내용 짠 뒤\n코드로 PPTX 파일 생성",
         "Manus · python-pptx · Marp", ORANGE),
    ]

    card_w = Inches(3.7)
    card_h = Inches(4.5)
    start_x = 0.95
    gap = 0.25

    for i, (letter, title, desc, examples, color) in enumerate(cards):
        x = Inches(start_x + i * (3.7 + gap))
        y = Inches(2.0)
        # 카드 배경
        add_round_rect(s, x, y, card_w, card_h, fill=WHITE, line=LIGHT_LINE)
        # 상단 색 바
        add_rect(s, x, y, card_w, Inches(0.5), fill=color)
        # 큰 알파벳
        add_text(s, x, Inches(2.55), card_w, Inches(1.0),
                 letter, size=56, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 제목
        add_text(s, Inches(start_x + i * (3.7 + gap) + 0.2), Inches(3.7),
                 Inches(3.3), Inches(0.5),
                 title, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # 설명
        add_text(s, Inches(start_x + i * (3.7 + gap) + 0.2), Inches(4.2),
                 Inches(3.3), Inches(1.2),
                 desc, size=12, color=TEXT, align=PP_ALIGN.CENTER)
        # 예시
        add_text(s, Inches(start_x + i * (3.7 + gap) + 0.2), Inches(5.7),
                 Inches(3.3), Inches(0.6),
                 examples, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


def slide_approach_detail(prs, page, total, *, letter, name, color,
                          tagline, pros, cons, recommended_for, example_tools):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, f"접근 {letter} — {name}",
           kicker=f"02-{letter}. 접근법 자세히", page=page, total=total)

    # 좌측 대형 알파벳 카드
    add_round_rect(s, Inches(0.95), Inches(2.0), Inches(3.5), Inches(4.5),
                   fill=color)
    add_text(s, Inches(0.95), Inches(2.0), Inches(3.5), Inches(2.0),
             letter, size=120, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.95), Inches(4.5), Inches(3.5), Inches(1.0),
             tagline, size=15, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.95), Inches(5.8), Inches(3.5), Inches(0.5),
             example_tools, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # 우측 텍스트 영역
    rx = Inches(5.0)
    rw = Inches(7.6)

    add_text(s, rx, Inches(2.0), rw, Inches(0.35),
             "장점", size=12, bold=True, color=ORANGE)
    for i, p in enumerate(pros):
        add_text(s, rx, Inches(2.4 + i * 0.45), rw, Inches(0.4),
                 f"• {p}", size=14, color=TEXT)

    cons_y = 2.4 + len(pros) * 0.45 + 0.3
    add_text(s, rx, Inches(cons_y), rw, Inches(0.35),
             "단점", size=12, bold=True, color=MUTED)
    for i, c in enumerate(cons):
        add_text(s, rx, Inches(cons_y + 0.4 + i * 0.45), rw, Inches(0.4),
                 f"• {c}", size=14, color=TEXT)

    rec_y = cons_y + 0.4 + len(cons) * 0.45 + 0.3
    add_round_rect(s, rx, Inches(rec_y), rw, Inches(0.7), fill=BG, line=LIGHT_LINE)
    add_text(s, Inches(5.2), Inches(rec_y), rw, Inches(0.7),
             f"이런 사람에게 추천:  {recommended_for}",
             size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


def slide_when_to_use(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "어떤 걸 언제 쓸까?",
           kicker="03. 상황별 가이드", page=page, total=total)

    rows = [
        ("30분 안에 초안이 필요한 내부 자료", "A — Gamma", BLUE),
        ("회사 표준 템플릿이 필수인 외부 발표", "B — PowerPoint + Copilot", NAVY),
        ("매주·매월 같은 형식의 정기 보고", "C — 코드 자동화", ORANGE),
        ("외부 IR / 세일즈 / 굵직한 임원 보고", "수동 + AI 보조", MUTED),
        ("학습 자료 · 강의 슬라이드", "A 또는 B", BLUE),
    ]

    # 헤더
    add_rect(s, Inches(0.95), Inches(2.0), Inches(7.5), Inches(0.5), fill=NAVY)
    add_text(s, Inches(1.1), Inches(2.0), Inches(7), Inches(0.5),
             "상황", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(8.45), Inches(2.0), Inches(4.15), Inches(0.5), fill=NAVY)
    add_text(s, Inches(8.6), Inches(2.0), Inches(4), Inches(0.5),
             "추천", size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    y = 2.55
    for situation, recommendation, color in rows:
        add_rect(s, Inches(0.95), Inches(y), Inches(11.65), Inches(0.65),
                 fill=WHITE, line=LIGHT_LINE)
        add_text(s, Inches(1.1), Inches(y), Inches(7.2), Inches(0.65),
                 situation, size=13, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        # 추천 컬러 점
        add_round_rect(s, Inches(8.6), Inches(y + 0.15), Inches(0.35),
                       Inches(0.35), fill=color)
        add_text(s, Inches(9.0), Inches(y), Inches(3.5), Inches(0.65),
                 recommendation, size=13, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += 0.7

    add_text(s, Inches(0.95), Inches(6.6), Inches(11.7), Inches(0.4),
             "핵심:  자료의 \"수명\"과 \"청자의 기대치\"가 갈래를 결정한다.",
             size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_quick_start(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "30분 안에 첫 덱 만들기  (Gamma)",
           kicker="04. 실전", page=page, total=total)

    steps = [
        ("1", "5분", "메시지 정리", "Claude/ChatGPT에 묻기:\n\"이 발표에서 청자가 가져갈\n핵심 3가지는?\""),
        ("2", "10분", "Gamma에 프롬프트", "제목·청자·톤·분량·\n포함할 내용을 한 번에 입력"),
        ("3", "10분", "카드 단위 수정", "표현·숫자만 손보기\n디자인은 손대지 말 것"),
        ("4", "5분", "PPTX 내보내기", "다운로드 → 발표용 PC에서\n폰트만 한 번 확인"),
    ]

    step_w = Inches(2.85)
    for i, (num, time, title, desc) in enumerate(steps):
        x = Inches(0.95 + i * 3.0)
        # 카드
        add_round_rect(s, x, Inches(2.0), step_w, Inches(4.7),
                       fill=WHITE, line=LIGHT_LINE)
        # 상단 시간 배지
        add_round_rect(s, Inches(0.95 + i * 3.0 + 0.3), Inches(2.2),
                       Inches(1.0), Inches(0.45), fill=ORANGE)
        add_text(s, Inches(0.95 + i * 3.0 + 0.3), Inches(2.2),
                 Inches(1.0), Inches(0.45),
                 time, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 단계 번호 (큰 숫자)
        add_text(s, x, Inches(2.85), step_w, Inches(1.2),
                 num, size=72, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 제목
        add_text(s, x, Inches(4.4), step_w, Inches(0.5),
                 title, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # 설명
        add_text(s, Inches(0.95 + i * 3.0 + 0.2), Inches(5.0),
                 Inches(2.55), Inches(1.5),
                 desc, size=11, color=TEXT, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.95), Inches(6.9), Inches(11.7), Inches(0.3),
             "총 30분 — 디자인 잡일 시간을 메시지 다듬는 시간으로 돌려놓는 게 목표",
             size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_prompt_tips(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "좋은 프롬프트 4요소",
           kicker="05. 프롬프트", page=page, total=total)

    elements = [
        ("청자", "Audience", "임원? 신입? 투자자?\n톤이 완전히 달라진다", BLUE),
        ("목적", "Purpose", "의사결정? 학습?\n소개?", NAVY),
        ("분량·구조", "Format", "\"8장, 표지·목차·결론 포함\"\n같이 구체적으로", ORANGE),
        ("금지사항", "Constraints", "이모지 금지, 한 장에 글머리\n기호 3개 이내 등", MUTED),
    ]

    for i, (kor, eng, desc, color) in enumerate(elements):
        row, col = i // 2, i % 2
        x = Inches(0.95 + col * 6.0)
        y = Inches(2.0 + row * 2.4)
        add_round_rect(s, x, y, Inches(5.7), Inches(2.1),
                       fill=WHITE, line=LIGHT_LINE)
        # 색 바
        add_rect(s, x, y, Inches(0.15), Inches(2.1), fill=color)
        # 영문 라벨
        add_text(s, Inches(0.95 + col * 6.0 + 0.4), Inches(2.0 + row * 2.4 + 0.2),
                 Inches(5), Inches(0.3),
                 eng.upper(), size=10, bold=True, color=color)
        # 한글 제목
        add_text(s, Inches(0.95 + col * 6.0 + 0.4), Inches(2.0 + row * 2.4 + 0.5),
                 Inches(5), Inches(0.7),
                 kor, size=24, bold=True, color=NAVY)
        # 설명
        add_text(s, Inches(0.95 + col * 6.0 + 0.4), Inches(2.0 + row * 2.4 + 1.25),
                 Inches(5), Inches(0.9),
                 desc, size=12, color=TEXT)

    add_text(s, Inches(0.95), Inches(6.85), Inches(11.7), Inches(0.35),
             "이 4가지가 빠지면 AI는 \"디지털 혁신\" 같은 일반론으로 슬라이드를 채운다.",
             size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_prompt_example(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "좋은 프롬프트 vs 나쁜 프롬프트",
           kicker="05. 프롬프트 예시", page=page, total=total)

    # 좌측: 나쁜 예
    add_round_rect(s, Inches(0.95), Inches(2.0), Inches(5.6), Inches(4.7),
                   fill=WHITE, line=MUTED)
    add_text(s, Inches(1.15), Inches(2.15), Inches(5.3), Inches(0.4),
             "✗  나쁜 예", size=14, bold=True, color=MUTED)
    add_rect(s, Inches(1.15), Inches(2.7), Inches(5.3), Inches(2),
             fill=BG, line=LIGHT_LINE)
    add_text(s, Inches(1.35), Inches(2.9), Inches(5.0), Inches(1.8),
             "분기 보고 PPT 만들어줘",
             size=14, color=TEXT)
    add_text(s, Inches(1.15), Inches(5.0), Inches(5.3), Inches(1.5),
             "→ AI는 \"디지털 혁신\", \"고객 중심\"\n   같은 일반론으로 채운다.\n→ 결과를 거의 못 쓴다.",
             size=12, color=MUTED)

    # 우측: 좋은 예
    add_round_rect(s, Inches(6.85), Inches(2.0), Inches(5.7), Inches(4.7),
                   fill=WHITE, line=ORANGE, line_width=2)
    add_text(s, Inches(7.05), Inches(2.15), Inches(5.3), Inches(0.4),
             "✓  좋은 예", size=14, bold=True, color=ORANGE)
    good_text = (
        "청자: 회사 임원 5명\n"
        "목적: 2분기 예산 추가 승인\n"
        "분량: 8장\n"
        "  (표지 + 현황 2 + 문제 2 + 제안 2 + 결론)\n"
        "톤: 데이터 중심, 감성 표현 배제\n"
        "금지: 이모지, 일반론적 표현"
    )
    add_rect(s, Inches(7.05), Inches(2.7), Inches(5.3), Inches(3),
             fill=BG, line=LIGHT_LINE)
    add_text(s, Inches(7.25), Inches(2.85), Inches(5.0), Inches(2.8),
             good_text, size=12, color=TEXT)
    add_text(s, Inches(7.05), Inches(5.9), Inches(5.3), Inches(0.7),
             "→ 청자·분량·톤·금지가 명확하면\n   AI 결과의 절반 이상이 그대로 쓸만해진다.",
             size=12, bold=True, color=ORANGE)


def slide_pitfalls(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "흔한 함정 4가지",
           kicker="06. 주의점", page=page, total=total)

    pitfalls = [
        ("AI가 만든 티",
         "모든 슬라이드가 비슷한 아이콘·색감 →\n임원·고객이 즉시 알아챈다",
         "표지·결론은 사람이 직접 디자인",
         ORANGE),
        ("환각 (Hallucination)",
         "AI가 수치·인용·통계를 그럴듯하게 지어낸다",
         "숫자 슬라이드는 원자료와 반드시 대조",
         NAVY),
        ("보안·정보 유출",
         "사내 기밀을 외부 AI 서비스에 그대로\n붙여넣으면 유출 사고",
         "사내 정책 확인. 가능하면 사내 Copilot 사용",
         BLUE),
        ("회사 템플릿 무시",
         "AI 네이티브 도구로 만든 자료는\n회사 마스터와 어긋남",
         "외부용·공식 자료는 처음부터 PowerPoint",
         MUTED),
    ]

    for i, (title, problem, solution, color) in enumerate(pitfalls):
        row, col = i // 2, i % 2
        x = Inches(0.95 + col * 6.0)
        y = Inches(2.0 + row * 2.4)
        # 카드
        add_round_rect(s, x, y, Inches(5.7), Inches(2.1),
                       fill=WHITE, line=LIGHT_LINE)
        # 좌측 컬러 액센트
        add_rect(s, x, y, Inches(0.15), Inches(2.1), fill=color)
        # 제목
        add_text(s, Inches(0.95 + col * 6.0 + 0.4), Inches(2.0 + row * 2.4 + 0.15),
                 Inches(5.0), Inches(0.5),
                 f"⚠  {title}", size=16, bold=True, color=NAVY)
        # 문제
        add_text(s, Inches(0.95 + col * 6.0 + 0.4), Inches(2.0 + row * 2.4 + 0.65),
                 Inches(5.0), Inches(0.8),
                 problem, size=11, color=TEXT)
        # 해결책
        add_text(s, Inches(0.95 + col * 6.0 + 0.4), Inches(2.0 + row * 2.4 + 1.55),
                 Inches(5.0), Inches(0.5),
                 f"→  {solution}", size=11, bold=True, color=color)


def slide_closing(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)

    # 장식
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5),
                              Inches(6), Inches(6))
    deco.fill.solid()
    deco.fill.fore_color.rgb = BLUE
    deco.line.fill.background()
    deco.shadow.inherit = False

    add_text(s, Inches(0.95), Inches(1.5), Inches(6), Inches(0.4),
             "정리", size=14, bold=True, color=ORANGE)
    add_text(s, Inches(0.95), Inches(2.0), Inches(12), Inches(2.5),
             "도구가 아니라\n\"어디까지 사람이 손대는가\"가 핵심.",
             size=40, bold=True, color=WHITE)

    add_text(s, Inches(0.95), Inches(4.8), Inches(12), Inches(0.4),
             "분담 규칙", size=13, bold=True, color=ORANGE)

    add_text(s, Inches(0.95), Inches(5.2), Inches(11.5), Inches(0.5),
             "잡일 (초안·디자인·아이콘 배치)",
             size=14, color=BG)
    add_text(s, Inches(0.95), Inches(5.5), Inches(11.5), Inches(0.5),
             "→  AI", size=14, bold=True, color=ORANGE)

    add_text(s, Inches(0.95), Inches(6.0), Inches(11.5), Inches(0.5),
             "판단 (메시지·숫자·표지·결론)",
             size=14, color=BG)
    add_text(s, Inches(0.95), Inches(6.3), Inches(11.5), Inches(0.5),
             "→  사람", size=14, bold=True, color=WHITE)

    add_text(s, Inches(0.95), Inches(7.05), Inches(12), Inches(0.3),
             "다음 시리즈 — AI 활용 · 회의록 편   |   2026.05.12",
             size=10, color=BG)


# ────────────────────────────────────────────────────────────────────
# 메인
# ────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    TOTAL = 13

    slide_cover(prs)                                # 1
    slide_agenda(prs, 2, TOTAL)                     # 2
    slide_problem(prs, 3, TOTAL)                    # 3
    slide_three_approaches(prs, 4, TOTAL)           # 4

    # 접근 A, B, C 상세
    slide_approach_detail(                          # 5
        prs, 5, TOTAL,
        letter="A", name="AI 네이티브 도구", color=BLUE,
        tagline="\"프롬프트 한 줄로\n전체 덱이 나온다\"",
        pros=[
            "디자인 고민 0 — 가장 빠른 초안",
            "한 번에 전체 덱 구조까지 생성",
            "무료 플랜으로도 PPTX 내보내기"
        ],
        cons=[
            "회사 표준 템플릿 강제 적용 어려움",
            "정밀 편집은 PowerPoint보다 약함",
        ],
        recommended_for="내부용 자료를 30분 안에 끝내고 싶은 사람",
        example_tools="Gamma · Tome · Beautiful.ai",
    )
    slide_approach_detail(                          # 6
        prs, 6, TOTAL,
        letter="B", name="기존 도구 + AI", color=NAVY,
        tagline="\"우리 회사 PPT\n그대로 자동 생성\"",
        pros=[
            "회사 마스터 슬라이드에 자연스럽게 융합",
            "보안·SSO 친화 (사내 처리)",
            "기존 PPT 편집 기능 그대로",
        ],
        cons=[
            "유료 라이선스 필요 (M365 Copilot 등)",
            "디자인 신선도는 떨어짐",
        ],
        recommended_for="회사 외부 발표·임원 보고가 잦은 사람",
        example_tools="MS Copilot · Gemini · Canva",
    )
    slide_approach_detail(                          # 7
        prs, 7, TOTAL,
        letter="C", name="LLM + 변환 코드", color=ORANGE,
        tagline="\"매주 같은 보고서를\n자동으로 찍어낸다\"",
        pros=[
            "정기 보고서 완전 자동화 가능",
            "데이터 소스와 직접 연결",
            "이 가이드 자료도 C 방식으로 생성됨",
        ],
        cons=[
            "초기 셋업에 코딩 필요",
            "디자인 정밀도는 손이 가는 만큼",
        ],
        recommended_for="반복되는 정기 보고를 자동화하고 싶은 사람",
        example_tools="Manus · python-pptx · Marp",
    )

    slide_when_to_use(prs, 8, TOTAL)                # 8
    slide_quick_start(prs, 9, TOTAL)                # 9
    slide_prompt_tips(prs, 10, TOTAL)               # 10
    slide_prompt_example(prs, 11, TOTAL)            # 11
    slide_pitfalls(prs, 12, TOTAL)                  # 12
    slide_closing(prs, 13, TOTAL)                   # 13

    out = Path(__file__).parent / "AI로-PPT-만들기-입문자가이드.pptx"
    prs.save(out)
    print(f"생성 완료: {out}")
    print(f"슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
